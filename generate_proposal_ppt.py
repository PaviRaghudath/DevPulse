"""Generate FileAnalyzer manager proposal presentation (Who/What/How/Impact/Scope)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1E, 0x6F, 0xBF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x55, 0x65, 0x77)
GREEN  = RGBColor(0x1A, 0x8C, 0x5E)
ORANGE = RGBColor(0xE8, 0x7D, 0x1E)
PURPLE = RGBColor(0x6A, 0x1B, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h,
       bold=False, italic=False, size=14, color=WHITE,
       align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.name   = "Calibri"
    return txb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.35), fill=NAVY)
    tb(slide, title, Inches(0.45), Inches(0.12), Inches(10), Inches(0.7),
       bold=True, size=28, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(0.45), Inches(0.78), Inches(11.5), Inches(0.45),
           size=14, color=RGBColor(0xAA, 0xC4, 0xE8))
    rect(slide, 0, Inches(1.35), SLIDE_W, Inches(0.06), fill=BLUE)


def footer(slide, text="FileAnalyzer  |  Project Proposal  |  2026"):
    rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), fill=NAVY)
    tb(slide, text, Inches(0.3), Inches(7.19), Inches(12), Inches(0.28),
       size=10, color=RGBColor(0x88, 0xA8, 0xCC), align=PP_ALIGN.CENTER)


def tag_badge(slide, text, l, t, color=BLUE):
    rect(slide, l, t, Inches(2.0), Inches(0.35), fill=color)
    tb(slide, text, l + Inches(0.08), t + Inches(0.04), Inches(1.88), Inches(0.28),
       bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
rect(sl, Inches(8.8), 0, Inches(4.53), SLIDE_H, fill=BLUE)
rect(sl, Inches(8.4), 0, Inches(0.5), SLIDE_H, fill=RGBColor(0x14, 0x55, 0x9A))

# badge
rect(sl, Inches(0.45), Inches(1.7), Inches(2.6), Inches(0.38), fill=BLUE)
tb(sl, "MANAGER PRESENTATION", Inches(0.5), Inches(1.73), Inches(2.55), Inches(0.32),
   size=10, bold=True, color=WHITE)

tb(sl, "FileAnalyzer", Inches(0.45), Inches(2.3), Inches(7.8), Inches(1.1),
   bold=True, size=52, color=WHITE)
tb(sl, "AI-Powered Log Monitoring, Document Analysis\n& Deployment Automation",
   Inches(0.45), Inches(3.35), Inches(7.8), Inches(0.85),
   size=20, color=RGBColor(0xAA, 0xC4, 0xE8))

for i, line in enumerate([
    "Automated error detection across remote servers",
    "AI document Q&A — any file, plain English answers",
    "One-click Git integration & CI/CD deployment",
]):
    tb(sl, f"›  {line}", Inches(0.55), Inches(4.35) + Inches(0.42)*i,
       Inches(7.6), Inches(0.38), size=14, color=RGBColor(0xCC, 0xDD, 0xF0))

# right panel info boxes
for i, (label, value) in enumerate([
    ("Who",    "Dev, DevOps, Support"),
    ("Status", "Production Ready"),
    ("Scope",  "Current + Future"),
]):
    y = Inches(1.4) + Inches(1.6)*i
    rect(sl, Inches(9.1), y, Inches(3.8), Inches(1.2), fill=RGBColor(0x0A, 0x12, 0x1F))
    tb(sl, label, Inches(9.22), y + Inches(0.1), Inches(3.5), Inches(0.38),
       size=11, color=RGBColor(0x88, 0xBB, 0xFF))
    tb(sl, value, Inches(9.22), y + Inches(0.45), Inches(3.5), Inches(0.55),
       bold=True, size=17, color=WHITE)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Overview (Who / What / How / Impact table)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Proposal Overview", "At a glance — the problem, solution, and value")
footer(sl)

rows = [
    (BLUE,   "WHO",    "Development Teams  ·  DevOps / Operations  ·  Support Teams"),
    (GREEN,  "WHAT",   "Log monitoring is manual and slow. Document review requires reading entire files. "
                       "Deployments require direct server access and SSH expertise. "
                       "Teams must switch context away from their IDE to investigate issues."),
    (ORANGE, "HOW",    "FileAnalyzer connects to remote servers via SSH, pulls logs on a schedule, "
                       "detects errors automatically, and sends HTML email alerts. "
                       "Any document (PDF, Word, CSV) can be queried in plain English using AI. "
                       "Git integration detects project type and triggers CI/CD or SSH build+deploy. "
                       "A REST API lets VS Code, IntelliJ, and Visual Studio connect directly."),
    (PURPLE, "IMPACT", "Faster error detection  ·  No manual log tailing  ·  "
                       "Zero context switching for developers  ·  Automated deployment pipeline  ·  "
                       "Instant AI answers from any document"),
]

row_h = [Inches(0.68), Inches(1.25), Inches(1.85), Inches(0.85)]
y = Inches(1.55)
for (color, label, text), h in zip(rows, row_h):
    rect(sl, Inches(0.35), y, Inches(1.55), h, fill=color)
    tb(sl, label, Inches(0.35), y + h/2 - Pt(10),
       Inches(1.55), Inches(0.45),
       bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, Inches(1.9), y, Inches(11.08), h, fill=LIGHT)
    rect(sl, Inches(1.9), y, Inches(0.06), h, fill=color)
    tb(sl, text, Inches(2.08), y + Inches(0.1),
       Inches(10.75), h - Inches(0.18),
       size=13, color=NAVY)
    y += h + Inches(0.08)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Problem Statement", "What teams are dealing with today")
footer(sl)

problems = [
    (ORANGE, "Manual Log Monitoring",
     "Developers SSH into production servers to tail log files when issues occur. "
     "There is no automated alerting — errors go unnoticed until a user reports them."),
    (ORANGE, "No Document Intelligence",
     "Reviewing a 200-page spec or contract requires reading the whole file. "
     "Teams have no way to ask plain English questions and get instant answers."),
    (ORANGE, "Complex Deployments",
     "Deploying a Spring Boot or React app requires knowing the server layout, "
     "running build commands over SSH, and copying artifacts manually."),
    (ORANGE, "Constant Context Switching",
     "Developers leave IntelliJ or VS Code to open a browser or terminal "
     "every time they need to check a log or look something up in a document."),
]

cw = Inches(5.95)
ch = Inches(2.18)
gx = Inches(0.43)
gy = Inches(0.25)
for i, (color, title, desc) in enumerate(problems):
    col_i = i % 2
    row_i = i // 2
    lx = Inches(0.35) + (cw + gx) * col_i
    ty = Inches(1.58) + (ch + gy) * row_i
    rect(sl, lx, ty, cw, ch, fill=LIGHT)
    rect(sl, lx, ty, cw, Inches(0.06), fill=color)
    rect(sl, lx, ty, Inches(0.07), ch, fill=color)
    tb(sl, title, lx + Inches(0.2), ty + Inches(0.1),
       cw - Inches(0.3), Inches(0.38),
       bold=True, size=15, color=NAVY)
    tb(sl, desc, lx + Inches(0.2), ty + Inches(0.52),
       cw - Inches(0.3), ch - Inches(0.6),
       size=12.5, color=GRAY)

rect(sl, 0, Inches(6.25), SLIDE_W, Inches(0.85), fill=NAVY)
tb(sl, "Result:  Revenue impact from missed errors, slow incident response, and developer frustration.",
   Inches(0.5), Inches(6.32), Inches(12.3), Inches(0.65),
   bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Proposed Solution
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Proposed Solution", "Four integrated capabilities — one AI agent")
footer(sl)

solutions = [
    (BLUE,   "Remote Log Monitoring",
     ["SSH into AWS EC2 servers on a schedule",
      "Incremental reads — only new log content fetched",
      "Auto-detect errors, warnings, and stack traces",
      "Send HTML email alerts when errors exceed threshold"]),
    (GREEN,  "AI Document Q&A",
     ["Upload PDF, Word, CSV, or log files",
      "Auto-detects document type and surfaces key findings",
      "Ask questions in plain English — get cited answers",
      "Supports OpenAI (GPT-4o) and Anthropic (Claude)"]),
    (ORANGE, "Git + CI/CD Integration",
     ["Connect a project's GitHub / GitLab URL",
      "Auto-detect type: Spring Boot JAR/WAR, React, Angular",
      "Trigger GitHub Actions, Jenkins, or GitLab CI pipelines",
      "Fall back to SSH build + deploy if no pipeline exists"]),
    (PURPLE, "IDE Hook (REST API)",
     ["Local REST API on localhost:8000",
      "VS Code, IntelliJ, Visual Studio can connect directly",
      "Send a log or document — get AI analysis inline",
      "No context switching — stay in your IDE"]),
]

cw = Inches(5.95)
ch = Inches(2.35)
gx = Inches(0.43)
gy = Inches(0.25)
for i, (color, title, bullets) in enumerate(solutions):
    col_i = i % 2
    row_i = i // 2
    lx = Inches(0.35) + (cw + gx) * col_i
    ty = Inches(1.55) + (ch + gy) * row_i
    rect(sl, lx, ty, cw, ch, fill=LIGHT)
    rect(sl, lx, ty, cw, Inches(0.48), fill=color)
    tb(sl, title, lx + Inches(0.15), ty + Inches(0.06),
       cw - Inches(0.25), Inches(0.38),
       bold=True, size=15, color=WHITE)
    for j, b in enumerate(bullets):
        tb(sl, f"›  {b}",
           lx + Inches(0.15),
           ty + Inches(0.58) + Inches(0.41)*j,
           cw - Inches(0.25), Inches(0.38),
           size=12.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Expected Impact
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Expected Impact", "Benefits per team")
footer(sl)

impacts = [
    (BLUE,   "Development Team",
     "Errors surfaced automatically — no manual log tailing or SSH sessions needed. "
     "Log analysis and document Q&A available directly inside VS Code or IntelliJ."),
    (GREEN,  "DevOps / Operations",
     "Scheduled monitoring of all Spring Boot and Apache servers from one dashboard. "
     "One-click deployment triggering CI/CD pipelines or SSH build+deploy."),
    (ORANGE, "Support Team",
     "Instant answers from design specs, contracts, and support runbooks in plain English. "
     "No need to read entire documents to find a single answer."),
    (PURPLE, "Business",
     "Faster incident response and fewer missed production errors. "
     "Reduced developer frustration and context switching overhead."),
]

row_h = Inches(1.18)
gap   = Inches(0.18)
y = Inches(1.58)
for color, team, desc in impacts:
    rect(sl, Inches(0.35), y, Inches(1.85), row_h, fill=color)
    tb(sl, team, Inches(0.35), y + row_h/2 - Pt(10),
       Inches(1.85), Inches(0.45),
       bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, Inches(2.2), y, Inches(10.78), row_h, fill=LIGHT)
    rect(sl, Inches(2.2), y, Inches(0.07), row_h, fill=color)
    tb(sl, desc, Inches(2.38), y + Inches(0.18),
       Inches(10.5), row_h - Inches(0.25),
       size=13.5, color=NAVY)
    y += row_h + gap

# KPI strip
rect(sl, 0, Inches(6.22), SLIDE_W, Inches(0.88), fill=NAVY)
kpis = [
    ("Automated",  "Error Detection"),
    ("0 SSH",      "Sessions to check logs"),
    ("4 IDEs",     "VS Code, IntelliJ, Visual Studio"),
    ("1-click",    "CI/CD Deployment"),
    ("Any file",   "PDF, Word, CSV, Log"),
]
for i, (val, label) in enumerate(kpis):
    lx = Inches(0.5) + Inches(2.55)*i
    tb(sl, val,   lx, Inches(6.26), Inches(2.4), Inches(0.45),
       bold=True, size=20, color=WHITE)
    tb(sl, label, lx, Inches(6.68), Inches(2.4), Inches(0.32),
       size=11, color=RGBColor(0xAA, 0xC4, 0xE8))


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Current Scope
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Current Scope", "What is built and ready today")
footer(sl)

scope_items = [
    (BLUE,   "1",
     "Document Analysis",
     "Upload any file (PDF, Word, CSV, .log, .txt). AI auto-detects document type, "
     "surfaces key findings, error counts, stack traces, and suggested questions."),
    (GREEN,  "2",
     "Remote Log Monitoring",
     "SSH into up to N AWS EC2 servers. Pull logs incrementally on a configurable schedule. "
     "Email alerts when errors exceed threshold — with top error table and stack trace count."),
    (ORANGE, "3",
     "Git Integration & Detection",
     "Connect a project's Git URL. Auto-detect type (Spring Boot JAR/WAR, React, Angular, Vue, "
     "Next.js) and build toolchain (Maven, Gradle, npm, yarn)."),
    (PURPLE, "4",
     "CI/CD Deployment",
     "Trigger GitHub Actions, Jenkins, or GitLab CI pipelines with one click. "
     "Fall back to SSH git pull → build → deploy for projects without a pipeline."),
    (RGBColor(0x00, 0x7A, 0x8A), "5",
     "IDE Hook — REST API",
     "POST /analyze and POST /ask endpoints on localhost:8000. "
     "IDE plugins (VS Code, IntelliJ, Visual Studio) send files and receive AI analysis without a browser."),
]

item_h = Inches(0.9)
gap    = Inches(0.16)
y = Inches(1.58)
for color, num, title, desc in scope_items:
    rect(sl, Inches(0.35), y, Inches(0.58), item_h, fill=color)
    tb(sl, num, Inches(0.35), y + Inches(0.2), Inches(0.58), Inches(0.45),
       bold=True, size=18, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.93), y, Inches(12.05), item_h, fill=LIGHT)
    rect(sl, Inches(0.93), y, Inches(0.06), item_h, fill=color)
    tb(sl, title, Inches(1.1), y + Inches(0.07),
       Inches(2.6), Inches(0.38),
       bold=True, size=13.5, color=color)
    tb(sl, desc, Inches(3.75), y + Inches(0.1),
       Inches(9.1), item_h - Inches(0.2),
       size=12.5, color=NAVY)
    y += item_h + gap


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Future Scope
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header(sl, "Future Scope", "Planned phases after initial rollout")
footer(sl)

future = [
    (BLUE,   "Slack / Microsoft Teams Notifications",
     "Send alert summaries directly to team channels in addition to email. "
     "Developers get notified where they already work."),
    (GREEN,  "Auto-Fix Suggestions",
     "When a known error pattern is detected in a log, suggest a fix or link to a runbook — "
     "directly in the IDE or email alert."),
    (ORANGE, "Multi-Server Dashboard",
     "Aggregate health status, error trends, and deployment history across all projects "
     "in a single view with charts and drill-down."),
    (PURPLE, "Published IDE Plugins",
     "Package and publish VS Code extension to the Marketplace and IntelliJ plugin to "
     "the JetBrains Plugin Repository — one-click install for developers."),
    (RGBColor(0x00, 0x7A, 0x8A), "Scheduled Report Generation",
     "Weekly PDF summaries of error trends, deployment history, and document activity "
     "per project — delivered by email automatically."),
]

cw = Inches(3.75)
ch = Inches(2.18)
gx = Inches(0.25)

# row 1: 3 cards
for i, (color, title, desc) in enumerate(future[:3]):
    lx = Inches(0.35) + (cw + gx)*i
    ty = Inches(1.58)
    rect(sl, lx, ty, cw, ch, fill=LIGHT)
    rect(sl, lx, ty, cw, Inches(0.06), fill=color)
    rect(sl, lx, ty, Inches(0.07), ch, fill=color)
    tb(sl, title, lx + Inches(0.2), ty + Inches(0.1),
       cw - Inches(0.3), Inches(0.5),
       bold=True, size=13.5, color=NAVY)
    tb(sl, desc, lx + Inches(0.2), ty + Inches(0.65),
       cw - Inches(0.3), ch - Inches(0.75),
       size=12, color=GRAY)

# row 2: 2 wider cards
cw2 = Inches(5.9)
for i, (color, title, desc) in enumerate(future[3:]):
    lx = Inches(0.35) + (cw2 + Inches(0.43))*i
    ty = Inches(4.0)
    rect(sl, lx, ty, cw2, ch, fill=LIGHT)
    rect(sl, lx, ty, cw2, Inches(0.06), fill=color)
    rect(sl, lx, ty, Inches(0.07), ch, fill=color)
    tb(sl, title, lx + Inches(0.2), ty + Inches(0.1),
       cw2 - Inches(0.3), Inches(0.45),
       bold=True, size=13.5, color=NAVY)
    tb(sl, desc, lx + Inches(0.2), ty + Inches(0.6),
       cw2 - Inches(0.3), ch - Inches(0.7),
       size=12, color=GRAY)

rect(sl, 0, Inches(6.38), SLIDE_W, Inches(0.72), fill=NAVY)
tb(sl, "Future phases will be scoped and prioritised based on team feedback after the initial rollout.",
   Inches(0.5), Inches(6.44), Inches(12.3), Inches(0.5),
   italic=True, size=13, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Approval Request
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
rect(sl, Inches(8.8), 0, Inches(4.53), SLIDE_H, fill=BLUE)
rect(sl, Inches(8.4), 0, Inches(0.5), SLIDE_H, fill=RGBColor(0x14, 0x55, 0x9A))
footer(sl, "FileAnalyzer  |  Requesting Approval")

tb(sl, "Approval Request",
   Inches(0.5), Inches(0.35), Inches(8), Inches(0.65),
   bold=True, size=30, color=WHITE)
rect(sl, Inches(0.5), Inches(1.0), Inches(3.2), Inches(0.05), fill=BLUE)

rect(sl, Inches(0.4), Inches(1.15), Inches(7.8), Inches(3.5), fill=RGBColor(0x10, 0x1E, 0x35))
tb(sl, "Requesting approval to proceed with current scope:",
   Inches(0.6), Inches(1.22), Inches(7.4), Inches(0.4),
   bold=True, size=15, color=BLUE)

approve_items = [
    "Document Analysis — AI Q&A on any file type",
    "Remote Log Monitoring — SSH + scheduled error detection + email alerts",
    "Git Integration — auto-detect Spring Boot, React, Angular projects",
    "CI/CD Deployment — GitHub Actions, Jenkins, GitLab CI + SSH fallback",
    "IDE Hook — REST API for VS Code, IntelliJ, Visual Studio",
]
for j, item in enumerate(approve_items):
    tb(sl, f"✓  {item}", Inches(0.62), Inches(1.72) + Inches(0.48)*j,
       Inches(7.4), Inches(0.42),
       size=13.5, color=WHITE)

# right side: summary stats
for i, (val, label) in enumerate([
    ("5",        "Core capabilities ready"),
    ("0",        "New infrastructure required"),
    ("Local",    "All AI runs on your machine"),
    ("~30 min",  "To deploy on AWS EC2"),
]):
    y = Inches(1.3) + Inches(1.4)*i
    rect(sl, Inches(9.2), y, Inches(3.7), Inches(1.15), fill=RGBColor(0x0A, 0x12, 0x1F))
    tb(sl, val,   Inches(9.35), y + Inches(0.1), Inches(3.4), Inches(0.55),
       bold=True, size=26, color=WHITE)
    tb(sl, label, Inches(9.35), y + Inches(0.62), Inches(3.4), Inches(0.42),
       size=12, color=RGBColor(0x88, 0xBB, 0xFF))

rect(sl, Inches(0.4), Inches(4.85), Inches(7.8), Inches(0.75), fill=BLUE)
tb(sl, "FileAnalyzer is built and production-ready. No additional development time needed to begin.",
   Inches(0.6), Inches(4.92), Inches(7.5), Inches(0.58),
   bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
out = r"E:\mkonnekt\FileAnalyzer\FileAnalyzer_Proposal.pptx"
prs.save(out)
print(f"Saved: {out}")

"""Generate FileAnalyzer manager presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # dark background / heading fills
BLUE    = RGBColor(0x1E, 0x6F, 0xBF)   # accent
LIGHT   = RGBColor(0xF0, 0xF4, 0xFA)   # card background
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x55, 0x65, 0x77)
GREEN   = RGBColor(0x1A, 0x8C, 0x5E)
ORANGE  = RGBColor(0xE8, 0x7D, 0x1E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def tb(slide, text, l, t, w, h,
        bold=False, italic=False, size=18, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.name   = font
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark navy top bar with title + optional subtitle."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), fill=NAVY)
    tb(slide, title,
       Inches(0.45), Inches(0.12), Inches(10), Inches(0.7),
       bold=True, size=28, color=WHITE)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.45), Inches(0.78), Inches(11), Inches(0.45),
           size=15, color=RGBColor(0xAA, 0xC4, 0xE8))


def accent_bar(slide, y=Inches(1.35), color=BLUE, h=Inches(0.06)):
    add_rect(slide, 0, y, SLIDE_W, h, fill=color)


def footer(slide, text="FileAnalyzer  |  Internal Presentation  |  2026"):
    add_rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), fill=NAVY)
    tb(slide, text, Inches(0.3), Inches(7.19), Inches(12), Inches(0.28),
       size=10, color=RGBColor(0x88, 0xA8, 0xCC), align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, fill=LIGHT, radius=False):
    s = add_rect(slide, l, t, w, h, fill=fill)
    return s


def bullet_block(slide, items, l, t, w, h, size=14, color=NAVY,
                 bullet="•", gap=Inches(0.32)):
    y = t
    for item in items:
        tb(slide, f"{bullet}  {item}", l, y, w, gap,
           size=size, color=color)
        y += gap


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title / Cover
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# full background
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

# accent diagonal stripe (visual) — just a blue rect rotated-ish using two rects
add_rect(sl, Inches(8.6), 0, Inches(4.73), SLIDE_H, fill=BLUE)
add_rect(sl, Inches(8.2), 0, Inches(0.55), SLIDE_H, fill=RGBColor(0x14, 0x55, 0x9A))

# tag
add_rect(sl, Inches(0.45), Inches(1.8), Inches(2.2), Inches(0.38), fill=BLUE)
tb(sl, "INTERNAL PRESENTATION", Inches(0.48), Inches(1.83), Inches(2.2), Inches(0.35),
   size=10, bold=True, color=WHITE)

# main title
tb(sl, "FileAnalyzer", Inches(0.45), Inches(2.4), Inches(7.8), Inches(1.1),
   bold=True, size=54, color=WHITE)
tb(sl, "RAG-Powered Document Q&A Agent",
   Inches(0.45), Inches(3.45), Inches(7.8), Inches(0.6),
   size=22, color=RGBColor(0xAA, 0xC4, 0xE8))

# sub-bullets
for i, line in enumerate([
    "Upload any document — get instant AI-powered answers",
    "Local embeddings · OpenAI & Anthropic · Scalable to AWS",
]):
    tb(sl, line, Inches(0.55), Inches(4.2) + Inches(0.42)*i,
       Inches(7.5), Inches(0.4), size=15,
       color=RGBColor(0xCC, 0xDD, 0xF0))

# right-panel labels
for i, (icon, label) in enumerate([
    ("PDF · DOCX · TXT · CSV",  "Supported Formats"),
    ("OpenAI  +  Anthropic",     "AI Providers"),
    ("AWS EC2  /  ECS Fargate",  "Deployment"),
]):
    y = Inches(1.5) + Inches(1.55)*i
    add_rect(sl, Inches(9.05), y, Inches(3.8), Inches(1.2),
             fill=RGBColor(0x0A, 0x12, 0x1F))
    tb(sl, label, Inches(9.15), y + Inches(0.1), Inches(3.6), Inches(0.38),
       size=11, color=RGBColor(0x88, 0xBB, 0xFF))
    tb(sl, icon,  Inches(9.15), y + Inches(0.45), Inches(3.6), Inches(0.55),
       bold=True, size=17, color=WHITE)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Who Can Use It & Why It's Useful
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(sl, "Who Can Use FileAnalyzer?",
           "Any team that needs fast, accurate answers from large documents")
accent_bar(sl)
footer(sl)

personas = [
    ("Analysts &\nResearchers",
     ["Query 200-page reports instantly",
      "Extract KPIs without manual reading",
      "Cross-reference multiple documents"]),
    ("Legal &\nCompliance",
     ["Search contracts for specific clauses",
      "Flag policy violations in bulk",
      "Summarise regulatory filings"]),
    ("Operations &\nSupport",
     ["Answer SOPs / runbook questions",
      "Onboard new staff with doc Q&A",
      "Reduce ticket resolution time"]),
    ("Data &\nEngineering",
     ["Query CSV datasets in plain English",
      "Explore schema documentation",
      "Audit logs and structured data"]),
]

cw = Inches(2.9)
ch = Inches(4.2)
gap = Inches(0.22)
for i, (title, bullets) in enumerate(personas):
    lx = Inches(0.35) + (cw + gap) * i
    ty = Inches(1.55)
    card(sl, lx, ty, cw, ch, fill=LIGHT)
    add_rect(sl, lx, ty, cw, Inches(0.06), fill=BLUE)
    tb(sl, title, lx + Inches(0.15), ty + Inches(0.1), cw - Inches(0.2), Inches(0.7),
       bold=True, size=15, color=NAVY)
    for j, b in enumerate(bullets):
        tb(sl, f"› {b}",
           lx + Inches(0.15),
           ty + Inches(0.85) + Inches(0.42)*j,
           cw - Inches(0.25), Inches(0.4),
           size=13, color=GRAY)

# bottom value-prop strip
add_rect(sl, 0, Inches(6.0), SLIDE_W, Inches(1.1), fill=NAVY)
for i, kpi in enumerate([
    ("500 MB",   "Max file size"),
    ("< 3 sec",  "Avg query latency"),
    ("~400 MB",  "Peak RAM usage"),
    ("Free",     "Local embedding (no API cost)"),
    ("6",        "Source excerpts per answer"),
]):
    lx = Inches(0.5) + Inches(2.56)*i
    tb(sl, kpi[0], lx, Inches(6.05), Inches(2.4), Inches(0.5),
       bold=True, size=22, color=WHITE)
    tb(sl, kpi[1], lx, Inches(6.52), Inches(2.4), Inches(0.35),
       size=11, color=RGBColor(0xAA, 0xC4, 0xE8))


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Pipeline / How It Works
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(sl, "How FileAnalyzer Works",
           "Two-phase pipeline: Ingestion (once) → Retrieval (per query)")
accent_bar(sl)
footer(sl)

steps = [
    (BLUE,                        "1. Parse",     "Streaming parser\n(PDF/DOCX/TXT/CSV)\nNever loads full\nfile into memory"),
    (RGBColor(0x14, 0x75, 0xCF),  "2. Chunk",     "512-char overlapping\nwindows (64-char\noverlap) preserves\ncross-boundary context"),
    (RGBColor(0x0E, 0x5E, 0xA8),  "3. Embed",     "sentence-transformers\nall-MiniLM-L6-v2\n(local, free)\nBatch size 64"),
    (RGBColor(0x08, 0x48, 0x82),  "4. Store",     "ChromaDB\npersistent vector\nstore on disk\n1 collection/doc"),
    (GREEN,                        "5. Retrieve",  "Top-6 similar\nchunks via cosine\nsimilarity\n≤ 12,000 chars ctx"),
    (RGBColor(0x0F, 0x6E, 0x49),  "6. Answer",    "LLMClient streams\ntokens via\nOpenAI or\nAnthropic API"),
]

bw = Inches(1.95)
bh = Inches(3.8)
gap = Inches(0.18)
for i, (col, title, desc) in enumerate(steps):
    lx = Inches(0.3) + (bw + gap)*i
    ty = Inches(1.6)
    add_rect(sl, lx, ty, bw, Inches(0.5), fill=col)
    tb(sl, title, lx, ty + Inches(0.06), bw, Inches(0.38),
       bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    card(sl, lx, ty + Inches(0.5), bw, bh - Inches(0.5), fill=LIGHT)
    tb(sl, desc, lx + Inches(0.12), ty + Inches(0.62), bw - Inches(0.2), bh - Inches(0.7),
       size=12, color=NAVY)
    # arrow
    if i < len(steps)-1:
        ax = lx + bw + Inches(0.03)
        tb(sl, "→", ax, ty + Inches(1.6), Inches(0.18), Inches(0.35),
           bold=True, size=18, color=BLUE)

# ingestion vs retrieval label
add_rect(sl, Inches(0.3),  Inches(5.55), Inches(7.64), Inches(0.3), fill=RGBColor(0xD0, 0xE4, 0xF7))
add_rect(sl, Inches(8.07), Inches(5.55), Inches(4.93), Inches(0.3), fill=RGBColor(0xC8, 0xEE, 0xDE))
tb(sl, "INGESTION PHASE  (run once per document)",
   Inches(0.35), Inches(5.56), Inches(7.5), Inches(0.28),
   size=11, color=NAVY, bold=True)
tb(sl, "RETRIEVAL PHASE  (per user query)",
   Inches(8.12), Inches(5.56), Inches(4.8), Inches(0.28),
   size=11, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — AI Models
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(sl, "AI Models & Providers",
           "Switch between OpenAI and Anthropic in the UI sidebar — no code changes needed")
accent_bar(sl)
footer(sl)

# Embedding model card (full width)
add_rect(sl, Inches(0.35), Inches(1.55), Inches(12.6), Inches(1.05), fill=NAVY)
tb(sl, "Embedding Model  —  all-MiniLM-L6-v2  (sentence-transformers)",
   Inches(0.55), Inches(1.6), Inches(9), Inches(0.45),
   bold=True, size=16, color=WHITE)
for i, item in enumerate([
    "Runs 100% locally — zero API cost",
    "~80 MB download, cached after first run",
    "384-dimension vectors · optimised for semantic similarity",
    "Works offline after initial download",
]):
    tb(sl, f"›  {item}",
       Inches(0.55) + Inches(3.2)*i, Inches(2.03),
       Inches(3.1), Inches(0.35),
       size=12, color=RGBColor(0xAA, 0xC4, 0xE8))

# OpenAI panel
add_rect(sl, Inches(0.35), Inches(2.78), Inches(6.0), Inches(3.8), fill=LIGHT)
add_rect(sl, Inches(0.35), Inches(2.78), Inches(6.0), Inches(0.06), fill=BLUE)
tb(sl, "OpenAI  (ChatGPT)", Inches(0.55), Inches(2.82), Inches(5.6), Inches(0.45),
   bold=True, size=17, color=NAVY)

oai_rows = [
    ("gpt-4o",       "Fast",      "Best",      "$0.005 / 1K tokens",   "Recommended for most queries"),
    ("gpt-4o-mini",  "Very fast", "Good",      "$0.00015 / 1K tokens", "High-volume / cost-sensitive"),
    ("gpt-4-turbo",  "Moderate",  "Excellent", "$0.01 / 1K tokens",    "Deep analytical tasks"),
]
headers = ["Model", "Speed", "Quality", "Cost", "Best For"]
col_xs  = [Inches(0.45), Inches(1.6), Inches(2.35), Inches(3.05), Inches(4.25)]
col_ws  = [Inches(1.1),  Inches(0.7), Inches(0.65), Inches(1.15), Inches(1.85)]

add_rect(sl, Inches(0.37), Inches(3.33), Inches(5.96), Inches(0.32), fill=BLUE)
for ci, (h, x, w) in enumerate(zip(headers, col_xs, col_ws)):
    tb(sl, h, x, Inches(3.35), w, Inches(0.28),
       bold=True, size=11, color=WHITE)
for ri, row in enumerate(oai_rows):
    ry = Inches(3.68) + Inches(0.37)*ri
    bg = LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, Inches(0.37), ry, Inches(5.96), Inches(0.36), fill=bg)
    for ci, (cell, x, w) in enumerate(zip(row, col_xs, col_ws)):
        tb(sl, cell, x, ry + Inches(0.04), w, Inches(0.3),
           size=11, color=NAVY)

# Anthropic panel
add_rect(sl, Inches(6.7), Inches(2.78), Inches(6.28), Inches(3.8), fill=LIGHT)
add_rect(sl, Inches(6.7), Inches(2.78), Inches(6.28), Inches(0.06), fill=ORANGE)
tb(sl, "Anthropic  (Claude)", Inches(6.9), Inches(2.82), Inches(5.8), Inches(0.45),
   bold=True, size=17, color=NAVY)

ant_rows = [
    ("claude-sonnet-4-6", "Fast",      "Excellent", "$0.003 / 1K tokens",   "Best balance of quality & cost"),
    ("claude-opus-4-6",   "Slower",    "Best",      "$0.015 / 1K tokens",   "Complex reasoning & analysis"),
    ("claude-haiku-4-5",  "Very fast", "Good",      "$0.00025 / 1K tokens", "High-volume / cost-sensitive"),
]
col_xs2 = [Inches(6.8), Inches(8.35), Inches(9.12), Inches(9.85), Inches(11.0)]

add_rect(sl, Inches(6.72), Inches(3.33), Inches(6.24), Inches(0.32), fill=ORANGE)
for ci, (h, x, w) in enumerate(zip(headers, col_xs2, col_ws)):
    tb(sl, h, x, Inches(3.35), w, Inches(0.28),
       bold=True, size=11, color=WHITE)
for ri, row in enumerate(ant_rows):
    ry = Inches(3.68) + Inches(0.37)*ri
    bg = LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, Inches(6.72), ry, Inches(6.24), Inches(0.36), fill=bg)
    for ci, (cell, x, w) in enumerate(zip(row, col_xs2, col_ws)):
        tb(sl, cell, x, ry + Inches(0.04), w, Inches(0.3),
           size=11, color=NAVY)

tb(sl, "Users enter their API key once in the sidebar — the app never stores it on disk.",
   Inches(0.45), Inches(6.7), Inches(12), Inches(0.35),
   size=13, italic=True, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Deployment Options
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(sl, "Deployment Options",
           "Three paths from laptop to production — pick the right fit for your team")
accent_bar(sl)
footer(sl)

deploy_cards = [
    (BLUE, "Local / Developer",
     ["streamlit run app.py",
      "Python 3.10+, pip install",
      "API key in .env or UI sidebar",
      "Data stored in data/vector_store/",
      "Port 8501 (configurable)"],
     "Best for: Single developer,\nprototyping, demos",
     "$0 / month"),
    (GREEN, "AWS EC2  (Recommended)",
     ["Ubuntu 22.04, t3.medium (min 4 GB RAM)",
      "systemd service — auto-starts on reboot",
      "nginx reverse proxy on port 80",
      "HTTPS via Certbot / Let's Encrypt",
      "One command deploy script included"],
     "Best for: Small team,\npersistent shared instance",
     "~$30 / month"),
    (ORANGE, "AWS ECS Fargate  (Scalable)",
     ["Docker container → ECR registry",
      "ECS Fargate — no server management",
      "EFS for persistent ChromaDB storage",
      "ALB load balancer for HA",
      "API keys in Secrets Manager"],
     "Best for: Teams, auto-scaling,\nproduction workloads",
     "~$12–35 / month"),
]

cw = Inches(3.9)
ch = Inches(4.9)
gap = Inches(0.27)
for i, (col, title, bullets, note, cost) in enumerate(deploy_cards):
    lx = Inches(0.3) + (cw + gap)*i
    ty = Inches(1.55)
    card(sl, lx, ty, cw, ch, fill=LIGHT)
    add_rect(sl, lx, ty, cw, Inches(0.55), fill=col)
    tb(sl, title, lx + Inches(0.15), ty + Inches(0.08), cw - Inches(0.2), Inches(0.42),
       bold=True, size=16, color=WHITE)
    for j, b in enumerate(bullets):
        tb(sl, f"›  {b}",
           lx + Inches(0.15),
           ty + Inches(0.68) + Inches(0.44)*j,
           cw - Inches(0.25), Inches(0.42),
           size=12.5, color=NAVY)
    # note box
    add_rect(sl, lx + Inches(0.12), ty + Inches(3.5), cw - Inches(0.24), Inches(0.72), fill=col)
    tb(sl, note,
       lx + Inches(0.18), ty + Inches(3.52), cw - Inches(0.3), Inches(0.68),
       size=11.5, color=WHITE, italic=True)
    # cost badge
    add_rect(sl, lx + cw - Inches(1.1), ty + Inches(0.08), Inches(1.0), Inches(0.34),
             fill=RGBColor(0xFF, 0xFF, 0xFF))
    tb(sl, cost,
       lx + cw - Inches(1.08), ty + Inches(0.1), Inches(0.95), Inches(0.3),
       bold=True, size=11, color=col, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Scaling Scenarios
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(sl, "Scaling FileAnalyzer",
           "Four proven patterns — apply as demand grows")
accent_bar(sl)
footer(sl)

scenarios = [
    (BLUE, "Scale UP\n(Bigger machine)",
     [
         "Upgrade EC2: t3.medium → t3.large → m6i.xlarge",
         "More RAM = larger embedding model (e.g. all-mpnet-base-v2)",
         "More CPUs = faster chunking & embedding throughput",
         "Zero code changes — just resize the instance",
     ],
     "Trigger: single user, large files (100–500 MB), slow ingestion"),
    (GREEN, "Scale OUT\n(More instances)",
     [
         "ECS Fargate service: set desired-count > 1",
         "ALB routes users to available containers",
         "Each container has its own ChromaDB on EFS (shared)",
         "Add Auto Scaling policy on CPU > 70%",
     ],
     "Trigger: multiple concurrent users, peak-hour slowdowns"),
    (ORANGE, "Scale DOCUMENTS\n(Larger knowledge base)",
     [
         "Replace ChromaDB with pgvector (PostgreSQL) or Pinecone",
         "Store all document collections in one shared DB",
         "Retriever queries across ALL docs simultaneously",
         "Enable multi-document Q&A from one question",
     ],
     "Trigger: >100 documents, org-wide knowledge base"),
    (RGBColor(0x6A, 0x1B, 0xAA), "Scale TEAMS\n(Multi-user / Auth)",
     [
         "Add Cognito / OAuth login layer in front of ALB",
         "User-scoped collections: each user sees only their docs",
         "Usage logging via CloudWatch — track API spend per user",
         "Rate limiting via ALB WAF rules",
     ],
     "Trigger: multiple departments, compliance requirements"),
]

cw = Inches(5.95)
ch = Inches(2.3)
gap_x = Inches(0.43)
gap_y = Inches(0.28)
for i, (col, title, bullets, note) in enumerate(scenarios):
    col_i = i % 2
    row_i = i // 2
    lx = Inches(0.3) + (cw + gap_x)*col_i
    ty = Inches(1.58) + (ch + gap_y)*row_i
    card(sl, lx, ty, cw, ch, fill=LIGHT)
    add_rect(sl, lx, ty, Inches(0.07), ch, fill=col)
    tb(sl, title, lx + Inches(0.18), ty + Inches(0.08), Inches(2.0), Inches(0.7),
       bold=True, size=14, color=col)
    for j, b in enumerate(bullets):
        tb(sl, f"›  {b}",
           lx + Inches(0.18),
           ty + Inches(0.78) + Inches(0.3)*j,
           cw - Inches(0.28), Inches(0.28),
           size=11.5, color=NAVY)
    tb(sl, f"When to apply:  {note}",
       lx + Inches(0.18), ty + ch - Inches(0.42), cw - Inches(0.28), Inches(0.38),
       size=11, italic=True, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Summary / Next Steps
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H*0.55, fill=RGBColor(0x0A, 0x12, 0x22))
footer(sl, "FileAnalyzer  |  Questions?")

tb(sl, "Summary & Next Steps",
   Inches(0.6), Inches(0.25), Inches(10), Inches(0.65),
   bold=True, size=30, color=WHITE)
add_rect(sl, Inches(0.6), Inches(0.9), Inches(2.8), Inches(0.05), fill=BLUE)

# left — summary bullets
summary = [
    "Answers questions on any PDF, DOCX, TXT, or CSV file",
    "Runs locally — documents never leave your environment",
    "Supports both OpenAI and Anthropic models",
    "Handles files up to 500 MB with constant ~400 MB RAM",
    "Ready to deploy on AWS EC2 or ECS Fargate",
    "Scales from single developer to org-wide knowledge base",
]
add_rect(sl, Inches(0.4), Inches(1.1), Inches(6.3), Inches(3.1), fill=RGBColor(0x10, 0x1E, 0x35))
tb(sl, "What you get", Inches(0.6), Inches(1.18), Inches(5.8), Inches(0.4),
   bold=True, size=15, color=BLUE)
for j, b in enumerate(summary):
    tb(sl, f"✓  {b}", Inches(0.62), Inches(1.62) + Inches(0.37)*j,
       Inches(5.9), Inches(0.35),
       size=13, color=WHITE)

# right — next steps
next_steps = [
    ("1", "Run locally",      "pip install -r requirements.txt\nstreamlit run app.py"),
    ("2", "Add API key",      "Enter OpenAI or Anthropic key\nin the UI sidebar"),
    ("3", "Upload & ask",     "Drag a document, click Analyze,\nstart asking questions"),
    ("4", "Deploy to AWS",    "Follow EC2 guide in README\n(~30 min, one script)"),
]
add_rect(sl, Inches(7.1), Inches(1.1), Inches(5.85), Inches(3.1), fill=RGBColor(0x10, 0x1E, 0x35))
tb(sl, "Getting started", Inches(7.3), Inches(1.18), Inches(5.4), Inches(0.4),
   bold=True, size=15, color=ORANGE)
for j, (num, step, detail) in enumerate(next_steps):
    ry = Inches(1.62) + Inches(0.72)*j
    add_rect(sl, Inches(7.22), ry, Inches(0.38), Inches(0.38), fill=ORANGE)
    tb(sl, num, Inches(7.22), ry, Inches(0.38), Inches(0.38),
       bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, step,   Inches(7.72), ry,               Inches(4.9), Inches(0.3),
       bold=True, size=13, color=WHITE)
    tb(sl, detail, Inches(7.72), ry + Inches(0.3), Inches(4.9), Inches(0.38),
       size=11, color=RGBColor(0xAA, 0xC4, 0xE8))

# bottom CTA
add_rect(sl, Inches(0.4), Inches(4.45), Inches(12.55), Inches(0.85), fill=BLUE)
tb(sl,
   "FileAnalyzer is production-ready.  Deploy today and give your team instant answers from any document.",
   Inches(0.6), Inches(4.53), Inches(12.2), Inches(0.65),
   bold=True, size=17, color=WHITE, align=PP_ALIGN.CENTER)

# tech tags
tags = ["Python 3.10+", "Streamlit", "ChromaDB", "sentence-transformers",
        "OpenAI SDK", "Anthropic SDK", "Docker", "AWS EC2 / ECS"]
for k, tag in enumerate(tags):
    add_rect(sl, Inches(0.4) + Inches(1.63)*k, Inches(5.55),
             Inches(1.55), Inches(0.34), fill=RGBColor(0x10, 0x1E, 0x35))
    tb(sl, tag, Inches(0.45) + Inches(1.63)*k, Inches(5.57),
       Inches(1.5), Inches(0.3),
       size=11, color=RGBColor(0x88, 0xBB, 0xFF), align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
out = r"E:\mkonnekt\FileAnalyzer\FileAnalyzer_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
